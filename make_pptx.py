from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PAPER = RGBColor(0xF1, 0xF1, 0xF1)
CREAM = RGBColor(0xDA, 0xCE, 0xBE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x16, 0x12)
MUTED = RGBColor(0x5E, 0x57, 0x4E)
ORANGE = RGBColor(0xDB, 0x8F, 0x2A)
ORANGE_DARK = RGBColor(0xB8, 0x74, 0x1C)
LINE = RGBColor(0xD4, 0xC8, 0xB8)

W = Inches(13.333)
H = Inches(7.5)
OUT = "Agentic_AI_Presentation.pptx"
TOTAL = 22
VERSION = "v1.6.0"


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=PAPER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def bar(slide, y=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Inches(0.12))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()


def star(slide, left, top, size):
    s = slide.shapes.add_shape(MSO_SHAPE.STAR_4_POINT, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = ORANGE
    s.line.fill.background()
    return s


def textbox(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def bullets(slide, l, t, w, h, items, size=20):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size=size, color=INK)
    return box


def card(slide, l, t, w, h, fill=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = LINE
    s.adjustments[0] = 0.08
    return s


def footer(slide, page):
    textbox(slide, Inches(0.6), Inches(7.12), Inches(9), Inches(0.3), f"Agentic AI  ·  {VERSION}  ·  Abhishek Mishra", 12, False, MUTED)
    textbox(slide, Inches(11.2), Inches(7.12), Inches(1.6), Inches(0.3), f"{page} / {TOTAL}", 12, False, MUTED, PP_ALIGN.RIGHT)


def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ORANGE
    stripe.line.fill.background()
    star(s, Inches(0.85), Inches(1.35), Inches(0.62))
    textbox(s, Inches(0.85), Inches(2.1), Inches(12), Inches(0.4), "PROJECT BRIEFING", 14, True, ORANGE_DARK)
    textbox(s, Inches(0.85), Inches(2.5), Inches(12), Inches(1.0), "Agentic AI", 54, True, INK, font="Georgia")
    textbox(
        s,
        Inches(0.85),
        Inches(3.55),
        Inches(11.2),
        Inches(1.0),
        "A local ReAct assistant that thinks, uses tools, reads files, runs code, and answers like a teammate.",
        22,
        False,
        MUTED,
    )
    textbox(s, Inches(0.85), Inches(5.2), Inches(11), Inches(0.35), "Features  ·  How it works  ·  Benefits  ·  Full project knowledge", 16, True, ORANGE_DARK)
    textbox(s, Inches(0.85), Inches(5.7), Inches(11), Inches(0.35), "Python  ·  FastAPI  ·  Next.js  ·  Groq    |    Built by Abhishek Mishra", 16, False, INK)
    textbox(s, Inches(0.85), Inches(6.15), Inches(11), Inches(0.35), f"{VERSION}  ·  http://127.0.0.1:3000  ·  Live on Render", 14, False, MUTED)


def content_slide(prs, title, items, page, size=20):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    bar(s)
    textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7), title, 32, True, INK, font="Georgia")
    bullets(s, Inches(0.8), Inches(1.35), Inches(11.6), Inches(5.5), items, size)
    footer(s, page)
    return s


def cards_slide(prs, title, cards, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    bar(s)
    textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.7), title, 32, True, INK, font="Georgia")
    n = len(cards)
    gap = Inches(0.28)
    left = Inches(0.7)
    usable = Inches(11.9)
    cw = int((usable - gap * (n - 1)) / n)
    for i, (h, body) in enumerate(cards):
        x = left + i * (cw + gap)
        card(s, x, Inches(1.4), cw, Inches(5.0), WHITE)
        textbox(s, x + Inches(0.28), Inches(1.65), cw - Inches(0.5), Inches(0.9), h, 18, True, ORANGE_DARK)
        textbox(s, x + Inches(0.28), Inches(2.55), cw - Inches(0.5), Inches(3.5), body, 15, False, INK)
    footer(s, page)


def compare_slide(prs, title, left_title, left_items, right_title, right_items, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    bar(s)
    textbox(s, Inches(0.7), Inches(0.32), Inches(12), Inches(0.55), title, 28, True, INK, font="Georgia")
    card(s, Inches(0.55), Inches(1.05), Inches(5.9), Inches(5.8), WHITE)
    card(s, Inches(6.85), Inches(1.05), Inches(5.9), Inches(5.8), WHITE)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(1.05), Inches(5.9), Inches(0.1))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ORANGE
    stripe.line.fill.background()
    textbox(s, Inches(0.85), Inches(1.25), Inches(5.4), Inches(0.5), left_title, 20, True, MUTED)
    textbox(s, Inches(7.15), Inches(1.25), Inches(5.4), Inches(0.5), right_title, 20, True, ORANGE_DARK)
    bullets(s, Inches(0.85), Inches(1.85), Inches(5.35), Inches(4.7), left_items, 16)
    bullets(s, Inches(7.15), Inches(1.85), Inches(5.35), Inches(4.7), right_items, 16)
    footer(s, page)
    return s


def rows_slide(prs, title, rows, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    bar(s)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.6), title, 30, True, INK, font="Georgia")
    for i, (h, body) in enumerate(rows):
        y = Inches(1.15) + i * Inches(1.12)
        card(s, Inches(0.7), y, Inches(11.9), Inches(1.0))
        textbox(s, Inches(0.95), y + Inches(0.12), Inches(11.4), Inches(0.32), h, 16, True, ORANGE_DARK)
        textbox(s, Inches(0.95), y + Inches(0.46), Inches(11.4), Inches(0.42), body, 15, False, INK)
    footer(s, page)


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    title_slide(prs)

    content_slide(
        prs,
        "Agenda",
        [
            "What Agentic AI is, and how it differs from a normal chatbot",
            "The problem it solves and the benefits for a student, shop, or team",
            "Architecture: Next.js UI + FastAPI agent + Groq LLM",
            "Every feature in v1.6.0 — chat, files, Teach, Projects, Notes, sandbox, artifacts",
            "How the ReAct loop thinks, calls tools, and answers",
            "How to run it locally, demo it, and what this project is not",
        ],
        2,
    )

    compare_slide(
        prs,
        "Chatbot vs Agentic AI",
        "Normal chatbot",
        [
            "User types → model replies with text",
            "Cannot search live prices or weather",
            "Cannot do exact math or run code",
            "Forgets your shop, files, and style next chat",
            "You copy-paste everything yourself",
            "Looks smart, but it cannot act",
        ],
        "This project — an agent",
        [
            "Thinks: answer now, or call a tool?",
            "Acts: search, weather, math, notes, Python",
            "Observes the tool result, then continues",
            "Reads PDFs, photos, and project files",
            "Remembers Teach notes and project instructions",
            "Shows the work in chat, Artifacts, and Notes",
        ],
        3,
    )

    content_slide(
        prs,
        "The idea",
        [
            "The model is not hidden inside LangGraph or CrewAI — the ReAct loop is our own Python.",
            "You can see Think → Act → Observe → Answer. That is the whole point of this learning project.",
            "The chat feels like a teammate: cream/gold UI, name greeting, voice, files, and a side panel for work.",
            "Tools stay behind the scenes. The user sees “Searching the web…” then a clear answer with sources.",
            "It runs on your machine (or one Docker host). The Groq key stays in .env — never shown in the UI.",
        ],
        4,
    )

    content_slide(
        prs,
        "Benefits",
        [
            "Faster answers with live facts — weather, gold, news — instead of guessing from training data.",
            "One place for a shop or class: standing instructions, PDFs, and chats live inside a Project.",
            "You can teach it your name, tone, and facts without training a new neural model.",
            "It can draft a webpage or graphic and open it in Artifacts — not just dump code in chat.",
            "Short Python runs in a locked sandbox, so numbers and scripts are checked, not invented.",
            "Hindi and Bhojpuri replies, voice in/out, so it works for people who prefer speaking.",
            "You learn real agent design: tools, memory, streaming, auth — not a black-box demo.",
        ],
        5,
        size=18,
    )

    cards_slide(
        prs,
        "Who it helps",
        [
            (
                "Student / builder",
                "See a real ReAct loop, add a tool, ship a UI. Better than only watching a framework tutorial.",
            ),
            (
                "Shop or small team",
                "Teach brand facts, drop a price list PDF into a Project, ask in Hindi, export the chat.",
            ),
            (
                "Anyone who asks daily",
                "Weather, math, GitHub lookup, “make me a landing page”, “save this as a note”.",
            ),
        ],
        6,
    )

    cards_slide(
        prs,
        "Architecture",
        [
            (
                "Frontend",
                "Next.js 15 · React · TypeScript\n\nChat, sidebar, Artifacts, Settings, Projects, Notes.\n\nhttp://127.0.0.1:3000\nProxies /api to Python so cookies and SSE streaming work.",
            ),
            (
                "Backend",
                "FastAPI + Uvicorn\n\nLogin, chat stream, extract files, Teach, notes APIs, tool install.\n\nhttp://127.0.0.1:7860",
            ),
            (
                "Agent core",
                "Custom ReAct Agent\n(not LangGraph)\n\nLLM + tools + memory.\nOptional crew: Researcher then Writer.\nSandbox for short Python.",
            ),
        ],
        7,
    )

    content_slide(
        prs,
        "Tech stack",
        [
            "LLM: Groq openai/gpt-oss-120b (OpenAI-compatible API). Also OpenAI, OpenRouter, or Ollama.",
            "Vision: Groq qwen/qwen3.6-27b — reads photos from gallery or camera.",
            "Python 3.10+  ·  FastAPI  ·  Uvicorn  ·  httpx  ·  itsdangerous  ·  pypdf  ·  python-multipart",
            "Next.js 15  ·  React 19  ·  TypeScript  ·  marked for chat markdown",
            "Search: DuckDuckGo (ddgs). Weather: wttr.in. Notes saved as .md on disk.",
            "UI: cream sidebar #dacebe, chat paper #f1f1f1, gold accent #db8f2a, Fraunces + Figtree.",
            "Deploy: one Docker image on Render or Hugging Face Spaces. Vercel alone cannot run the Python API.",
        ],
        8,
        size=18,
    )

    content_slide(
        prs,
        "Feature map — what v1.6.0 can do",
        [
            "Chat: markdown, Recents search, edit a sent message, Retry, Export chat as .md",
            "Voice: microphone in, speaker out, English / Hindi / Bhojpuri",
            "Files: PDF, Word, txt, md, csv, json — plus photos and camera",
            "Teach: instructions, facts, and up to 5 training files on every chat (memory, not GPU training)",
            "Projects: name, standing instructions, files, and chats grouped together",
            "Notes viewer: agent saves markdown; you open, download, or delete it in the app",
            "Artifacts: HTML / SVG / markdown / code opens in a side pane",
            "Code sandbox: short Python, no network, no files, 6 second limit",
        ],
        9,
        size=18,
    )

    rows_slide(
        prs,
        "Chat experience",
        [
            ("Login & greeting", "Name + email, or Google / GitHub. Empty state: “Ready to leap, {name}?”"),
            ("Sidebar", "New leap, project picker, Single agent vs Researcher + Writer, Recents search, profile menu"),
            ("Composer", "+ attach (photo, camera, file), microphone, Enter to send (optional in Settings)"),
            ("After a reply", "Copy, speak, Retry last answer, or edit your bubble and Save & retry"),
            ("Language", "English, Hindi, or Bhojpuri — the agent is told to reply in that language"),
        ],
        10,
    )

    cards_slide(
        prs,
        "Files, photos, Teach",
        [
            (
                "Drop a file",
                "PDF, .docx, txt, md, csv, json. Text is extracted and wrapped as ---file:name--- so the agent quotes the real document.",
            ),
            (
                "Photos",
                "Up to 4 images, gallery or camera. Vision model answers from what is in the picture — it should not invent details.",
            ),
            (
                "Teach (Settings)",
                "How it should behave + facts to remember + training files (max 5). Injected on every chat. This is memory, not a new neural net.",
            ),
        ],
        11,
    )

    cards_slide(
        prs,
        "Projects, Notes, Artifacts",
        [
            (
                "Projects",
                "Like a folder: name, standing instructions, up to 8 files. Recents filter to that project. Context is sent with every message.",
            ),
            (
                "Notes",
                "Ask it to save a markdown note. Open Notes from the profile menu to read, download, or delete workspace_notes/*.md.",
            ),
            (
                "Artifacts",
                "Ask for a webpage, poster, or SVG. The agent emits an artifact tag. A pane on the right previews the file.",
            ),
        ],
        12,
    )

    content_slide(
        prs,
        "Code sandbox",
        [
            "Tool: code_run — core utility, always installed.",
            "Runs short Python in a child process (isolated interpreter, 6 second timeout).",
            "Allowed: math, json, re, datetime, statistics, random, and similar stdlib. print() is required for output.",
            "Blocked: os, subprocess, open(), eval/exec, network, installs, arbitrary files.",
            "Use it for exact numbers, small scripts, and checks — not as a full IDE or Claude Code clone.",
        ],
        13,
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s)
    bar(s)
    textbox(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.6), "Tools", 32, True, INK, font="Georgia")
    groups = [
        ("Core (ready by default)", "Web Search  ·  Wikipedia  ·  Weather  ·  Calculator  ·  Clock  ·  Notes  ·  Code sandbox"),
        ("Install when you need them", "GitHub lookup  ·  Dice  ·  Unit Convert  ·  Text Stats  ·  UUID  ·  Random Pick"),
        ("How they feel", "Settings → Tools: Install / Uninstall. The agent only calls what is installed. Status shows “Searching the web…”"),
    ]
    for i, (h, body) in enumerate(groups):
        y = Inches(1.25) + i * Inches(1.7)
        card(s, Inches(0.7), y, Inches(11.9), Inches(1.5))
        textbox(s, Inches(1.0), y + Inches(0.22), Inches(11.3), Inches(0.45), h, 20, True, ORANGE_DARK)
        textbox(s, Inches(1.0), y + Inches(0.7), Inches(11.3), Inches(0.6), body, 18, False, INK)
    footer(s, 14)

    content_slide(
        prs,
        "Two ways to work",
        [
            "Single agent — one assistant for quick questions, weather, prices, files, and code.",
            "Researcher + Writer — first agent gathers facts (search, wiki, weather, GitHub).",
            "Writer turns the brief into a clear answer and may use calculator, clock, notes, or sandbox.",
            "Switch from the sidebar dropdown before you send. Crew is better for long research.",
            "Show thinking: live status while tools run. Max steps default 8. Temperature default 0.2.",
        ],
        15,
    )

    content_slide(
        prs,
        "How the agent thinks (ReAct)",
        [
            "Think — the LLM decides: greet, answer, or call a tool.",
            "Act — it calls a tool with clean JSON, e.g. {\"query\": \"today gold price India\"}.",
            "Observe — the tool result is added to memory for the next step.",
            "Repeat — up to max steps, then write a human answer (not JSON, not a raw dump).",
            "Guardrails: repair bad Groq tool calls, cite sources, do not invent URLs or numbers.",
            "Extras in the prompt: language, Teach block, project files, attached file text, photos.",
        ],
        16,
    )

    content_slide(
        prs,
        "Login and security",
        [
            "Local login: name + email, signed httpOnly cookie session.",
            "Google and GitHub OAuth. Callback must use http://127.0.0.1:3000 (not localhost) for Google.",
            "Chat, tools, notes, extract, and settings APIs require a logged-in session.",
            "Groq API key lives in .env / server env. The UI never prints it back.",
            "Sandbox has no network and a short timeout. Teach is instruction memory, not model weights.",
            "Do not commit .env. Notes folder workspace_notes/ is gitignored.",
        ],
        17,
    )

    content_slide(
        prs,
        "How to run (Windows)",
        [
            "Terminal 1:  .\\.venv\\Scripts\\python.exe app.py     → API on port 7860",
            "Terminal 2:  cd frontend   then   npm run dev      → UI on port 3000",
            "Open http://127.0.0.1:3000  — if you see Connection Failed, these two processes are not running.",
            "Sign in, paste a Groq key if asked (console.groq.com/keys), then chat.",
            "CLI also works: python main.py  or  python main.py --crew -q \"...\"",
            "Live: Docker on Render (agentic-ai-d0cx.onrender.com) or Hugging Face Spaces.",
        ],
        18,
        size=18,
    )

    content_slide(
        prs,
        "Demo script",
        [
            "1. “Hello” — personal greeting, no tools.",
            "2. “Delhi weather” then “today gold price in India”.",
            "3. Attach a PDF: “summarize this file”.",
            "4. Settings → Teach: add your name and shop facts, Save training, new chat, ask “what do you know about me?”",
            "5. + Project, add a file, ask a question that needs that file.",
            "6. “Make a simple landing page for my shop” — Artifacts pane opens.",
            "7. “Run Python: print(sum(range(1, 11)))”. Then “save this as a note” and open Notes.",
            "8. Mic to speak, speaker to hear, Retry, Export chat.",
        ],
        19,
        size=18,
    )

    content_slide(
        prs,
        "Project structure",
        [
            "agentic_ai/agent.py — ReAct loop, Teach block, vision, artifacts instruction",
            "agentic_ai/tools.py — catalog + handlers (search, weather, notes, …)",
            "agentic_ai/sandbox.py — locked-down code_run",
            "agentic_ai/files.py — PDF / docx / text extract    prefs.py — settings.json + CORE_TOOLS",
            "agentic_ai/crew.py — Researcher + Writer    auth.py — sessions and OAuth",
            "app.py — FastAPI: /api/chat (SSE), /api/extract, /api/teach/*, /api/notes",
            "frontend/ — Next.js AppClient, ArtifactPane, Projects, Notes overlay",
            "main.py — CLI    Dockerfile — one-container deploy",
        ],
        20,
        size=17,
    )

    compare_slide(
        prs,
        "Honest scope — what it is not",
        "This app does",
        [
            "ReAct agent with real tools you can read in Python",
            "Chat UI with files, voice, projects, notes, artifacts",
            "Memory via Teach + project files (prompt context)",
            "Short sandboxed Python",
            "Local or one Docker host",
        ],
        "It does not (on purpose)",
        [
            "Clone Claude Code or Computer Use",
            "Fine-tune / train a new neural model on GPU",
            "Give the model your whole OS or browser",
            "Connect Gmail / Drive / Calendar (no extra OAuth yet)",
            "Replace LangGraph for huge production graphs",
        ],
        21,
    )

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, CREAM)
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ORANGE
    stripe.line.fill.background()
    star(s, Inches(0.85), Inches(1.7), Inches(0.62))
    textbox(s, Inches(0.85), Inches(2.5), Inches(12), Inches(0.9), "Thank you", 48, True, INK, font="Georgia")
    textbox(
        s,
        Inches(0.85),
        Inches(3.45),
        Inches(11.5),
        Inches(0.7),
        "Built and developed by Abhishek Mishra",
        22,
        True,
        ORANGE_DARK,
    )
    textbox(
        s,
        Inches(0.85),
        Inches(4.2),
        Inches(11.5),
        Inches(0.7),
        "Questions, ideas, or a live demo — open the app and try the script.",
        18,
        False,
        MUTED,
    )
    textbox(s, Inches(0.85), Inches(5.3), Inches(11), Inches(0.35), "http://127.0.0.1:3000", 18, True, INK)
    textbox(s, Inches(0.85), Inches(5.75), Inches(11), Inches(0.35), "Live: https://agentic-ai-d0cx.onrender.com", 16, False, MUTED)

    saved = []
    for path in (OUT, "Agentic_AI_Full_Briefing_v1.6.pptx"):
        try:
            prs.save(path)
            saved.append(path)
        except PermissionError:
            print("locked (close PowerPoint):", path)
    print("saved:", ", ".join(saved) if saved else "none")


if __name__ == "__main__":
    main()
